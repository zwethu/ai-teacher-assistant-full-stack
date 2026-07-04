"""Validation, processing, and persistence for conversation-scoped attachments."""

from __future__ import annotations

import csv
import hashlib
import io
import json
import logging
import math
import os
import uuid
import zipfile
import re
from datetime import datetime, timedelta, timezone
from pathlib import PurePath
from typing import Any

from google.cloud import firestore
from google.cloud.firestore import SERVER_TIMESTAMP

from entity.ChatAttachment import ChatAttachment
from services.attachment_constants import (
    ALLOWED_CONTENT_TYPES,
    CHARS_PER_TOKEN,
    DOCUMENT_CONTENT_TYPES,
    EXTENSION_CONTENT_TYPES,
    IMAGE_CONTENT_TYPES,
    IMAGE_TOKEN_ESTIMATE,
    MAX_CHAT_ATTACHMENT_BYTES,
    MAX_DOCUMENT_BYTES,
    MAX_EXTRACTED_PREVIEW_CHARS,
    MAX_FULL_EXTRACT_CHARS,
    MAX_IMAGE_BYTES,
    NATIVE_READABLE_CONTENT_TYPES,
    THUMBNAIL_MAX_SIZE,
    TOKENS_PER_PDF_PAGE,
    get_chat_attachment_retention_days,
)
from utils.firestore_client import get_firestore
from utils.gcs import (
    chat_attachment_blob_path,
    chat_attachment_thumbnail_path,
    safe_file_name,
    signed_read_url,
    upload_bytes,
    delete_blob,
)
from services.document_extraction import extract_document

logger = logging.getLogger(__name__)
ATTACHMENTS_SUBCOLLECTION = "attachments"


class AttachmentValidationError(ValueError):
    pass


def _chat_ref(batch_id: str, chat_id: str):
    return get_firestore().collection("batches").document(batch_id).collection("chats").document(chat_id)


def attachment_ref(batch_id: str, chat_id: str, attachment_id: str):
    return _chat_ref(batch_id, chat_id).collection(ATTACHMENTS_SUBCOLLECTION).document(attachment_id)


def _iso(value: Any) -> str | None:
    return value.isoformat() if hasattr(value, "isoformat") else (str(value) if value else None)


def attachment_to_model(doc_id: str, data: dict[str, Any]) -> ChatAttachment:
    kind = str(data.get("attachment_kind") or "other")
    status = str(data.get("status") or "processing")
    return ChatAttachment(
        attachment_id=doc_id,
        batch_id=str(data.get("batch_id") or ""),
        chat_id=str(data.get("chat_id") or ""),
        message_id=str(data["message_id"]) if data.get("message_id") else None,
        lecturer_id=str(data.get("lecturer_id") or ""),
        file_name=str(data.get("file_name") or ""),
        file_title=str(data.get("file_title") or data.get("file_name") or ""),
        content_type=str(data.get("content_type") or "application/octet-stream"),
        size_bytes=int(data.get("size_bytes") or 0),
        scope="chat",
        attachment_kind=kind if kind in {"document", "image", "other"} else "other",  # type: ignore[arg-type]
        status=status if status in {"processing", "ready", "failed"} else "processing",  # type: ignore[arg-type]
        content_sha256=str(data.get("content_sha256") or ""),
        token_estimate=int(data.get("token_estimate") or 0),
        parse_status=str(data.get("parse_status") or "pending"),  # type: ignore[arg-type]
        vision_status=str(data.get("vision_status") or "skipped"),  # type: ignore[arg-type]
        extracted_text_preview=str(data.get("extracted_text_preview") or ""),
        vision_summary=str(data.get("vision_summary") or ""),
        ocr_text=str(data.get("ocr_text") or ""),
        vision_error=str(data.get("vision_error") or ""),
        vision_source=str(data.get("vision_source") or "none"),  # type: ignore[arg-type]
        expires_at=_iso(data.get("expires_at")),
        promoted_file_id=None,
        promotion_allowed=False,
        thumbnail_available=bool(data.get("thumbnail_gcs_path")),
        rag_status=str(data.get("rag_status") or "skipped"),  # type: ignore[arg-type]
        chunk_status=str(data.get("chunk_status") or "skipped"),  # type: ignore[arg-type]
        embedding_status=str(data.get("embedding_status") or "skipped"),  # type: ignore[arg-type]
        semantic_search_ready=bool(data.get("semantic_search_ready", False)),
        chunk_count=int(data.get("chunk_count") or 0),
        indexed_chars=int(data.get("indexed_chars") or 0),
        ocr_status=str(data.get("ocr_status") or "not_needed"),  # type: ignore[arg-type]
        rag_updated_at=_iso(data.get("rag_updated_at")),
        created_at=_iso(data.get("created_at")),
        updated_at=_iso(data.get("updated_at")),
    )


def _normalize_content_type(content_type: str, extension: str) -> str:
    content_type = content_type.lower().split(";", 1)[0].strip()
    if extension == ".csv" and content_type in {"application/csv", "text/plain"}:
        return "text/csv"
    if extension == ".json" and content_type == "text/json":
        return "application/json"
    if extension in {".md", ".markdown"} and content_type == "text/plain":
        return "text/markdown"
    if extension in {".heic", ".heif"} and content_type in {"image/heic", "image/heif"}:
        return "image/heic" if extension == ".heic" else "image/heif"
    return content_type


def validate_attachment(file_name: str, content_type: str, data: bytes) -> tuple[str, str]:
    clean_name = safe_file_name(file_name)
    extension = PurePath(clean_name).suffix.lower()
    normalized = _normalize_content_type(content_type, extension)
    expected = EXTENSION_CONTENT_TYPES.get(extension)
    if not expected or normalized not in ALLOWED_CONTENT_TYPES or content_type.lower().split(";", 1)[0] not in expected:
        raise AttachmentValidationError("File extension and content type do not match or are unsupported.")
    limit = MAX_IMAGE_BYTES if normalized in IMAGE_CONTENT_TYPES else MAX_DOCUMENT_BYTES
    if not data:
        raise AttachmentValidationError("The uploaded file is empty.")
    if len(data) > limit:
        raise AttachmentValidationError(f"File exceeds the {limit // (1024 * 1024)} MB limit.")

    if normalized == "application/pdf" and not data.startswith(b"%PDF-"):
        raise AttachmentValidationError("Invalid PDF signature.")
    if normalized.startswith("application/vnd.openxmlformats"):
        _validate_ooxml(data, normalized)
    elif normalized in IMAGE_CONTENT_TYPES:
        _validate_image(data, normalized)
    elif normalized in {"text/plain", "text/markdown", "text/csv"}:
        if b"\x00" in data[:8192]:
            raise AttachmentValidationError("Text files cannot contain binary data.")
        try:
            data[:65536].decode("utf-8-sig")
        except UnicodeDecodeError as exc:
            raise AttachmentValidationError("Text files must be UTF-8 encoded.") from exc
    return clean_name, normalized


def validate_batch_document(file_name: str, content_type: str, data: bytes, max_bytes: int) -> tuple[str, str]:
    """Apply the same signature/OOXML policy to the batch-specific document set."""
    clean_name = safe_file_name(file_name)
    extension = PurePath(clean_name).suffix.lower()
    raw_type = content_type.lower().split(";", 1)[0].strip()
    normalized = _normalize_content_type(raw_type, extension)
    allowed = DOCUMENT_CONTENT_TYPES | {"application/json"}
    expected = EXTENSION_CONTENT_TYPES.get(extension)
    if not expected or normalized not in allowed or raw_type not in expected:
        raise AttachmentValidationError("File extension and content type do not match or are unsupported.")
    if not data:
        raise AttachmentValidationError("The uploaded file is empty.")
    if len(data) > max_bytes:
        raise AttachmentValidationError(f"File exceeds the {max_bytes // (1024 * 1024)} MB limit.")
    if normalized == "application/pdf" and not data.startswith(b"%PDF-"):
        raise AttachmentValidationError("Invalid PDF signature.")
    if normalized.startswith("application/vnd.openxmlformats"):
        _validate_ooxml(data, normalized)
    elif normalized in {"text/plain", "text/markdown", "text/csv", "application/json"}:
        if b"\x00" in data[:8192]:
            raise AttachmentValidationError("Text files cannot contain binary data.")
        try:
            decoded = data.decode("utf-8-sig")
            if normalized == "application/json":
                json.loads(decoded)
        except (UnicodeDecodeError, json.JSONDecodeError) as exc:
            raise AttachmentValidationError("Text and JSON files must be valid UTF-8 content.") from exc
    return clean_name, normalized


def _validate_ooxml(data: bytes, content_type: str) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            infos = archive.infolist()
            if len(infos) > 10_000 or sum(item.file_size for item in infos) > 200 * 1024 * 1024:
                raise AttachmentValidationError("The Office file expands beyond safe processing limits.")
            if any(item.flag_bits & 0x1 for item in infos):
                raise AttachmentValidationError("Encrypted Office files are not supported.")
            names = set(archive.namelist())
            if any(name.lower().endswith("vbaproject.bin") for name in names):
                raise AttachmentValidationError("Macro-enabled Office files are not allowed.")
            required = "word/document.xml" if "wordprocessingml" in content_type else "ppt/presentation.xml"
            if required not in names or "[Content_Types].xml" not in names:
                raise AttachmentValidationError("The Office file is malformed or has the wrong type.")
    except zipfile.BadZipFile as exc:
        raise AttachmentValidationError("The Office file is malformed.") from exc


def _register_heif() -> None:
    try:
        from pillow_heif import register_heif_opener
        register_heif_opener()
    except ImportError:
        pass


def _validate_image(data: bytes, content_type: str) -> None:
    _register_heif()
    from PIL import Image, UnidentifiedImageError
    expected = {
        "image/png": {"PNG"}, "image/jpeg": {"JPEG"}, "image/webp": {"WEBP"},
        "image/heic": {"HEIF", "HEIC"}, "image/heif": {"HEIF", "HEIC"},
    }[content_type]
    try:
        with Image.open(io.BytesIO(data)) as image:
            if str(image.format or "").upper() not in expected:
                raise AttachmentValidationError("Image signature does not match its content type.")
            image.verify()
    except (UnidentifiedImageError, OSError) as exc:
        raise AttachmentValidationError("The image is malformed or unsupported.") from exc


def _extract_document(data: bytes, content_type: str) -> str:
    limit = MAX_EXTRACTED_PREVIEW_CHARS
    if content_type == "application/pdf":
        from pypdf import PdfReader
        parts: list[str] = []
        size = 0
        for page in PdfReader(io.BytesIO(data)).pages:
            text = page.extract_text() or ""
            parts.append(text)
            size += len(text)
            if size >= limit:
                break
        return "\n\n".join(parts)[:limit]
    if "wordprocessingml" in content_type:
        from docx import Document
        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        size = 0
        for paragraph in doc.paragraphs:
            parts.append(paragraph.text)
            size += len(paragraph.text)
            if size >= limit:
                break
        return "\n".join(parts)[:limit]
    if "presentationml" in content_type:
        from pptx import Presentation
        presentation = Presentation(io.BytesIO(data))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
                    if sum(len(line) for line in lines) >= limit:
                        return "\n".join(lines)[:limit]
        return "\n".join(lines)[:limit]
    text = data[: max(limit * 4, limit)].decode("utf-8-sig", errors="ignore")
    if content_type == "text/csv":
        rows = csv.reader(io.StringIO(text))
        output: list[str] = []
        size = 0
        for row in rows:
            line = " | ".join(cell.strip() for cell in row)
            output.append(line)
            size += len(line)
            if size >= limit:
                break
        return "\n".join(output)[:limit]
    return text[:limit]


def _thumbnail_bytes(data: bytes) -> bytes:
    _register_heif()
    from PIL import Image, ImageOps
    with Image.open(io.BytesIO(data)) as image:
        image = ImageOps.exif_transpose(image).convert("RGB")
        image.thumbnail(THUMBNAIL_MAX_SIZE)
        output = io.BytesIO()
        image.save(output, format="JPEG", quality=84, optimize=True)
        return output.getvalue()


def _vision_context(data: bytes, gcs_path: str, content_type: str) -> tuple[str, str, str, str, str]:
    model = (os.getenv("ATTACHMENT_VISION_MODEL") or "").strip()
    project = (os.getenv("GOOGLE_CLOUD_PROJECT") or "").strip()
    if not model or not project:
        return "skipped", "", "", "Vision processing is not configured.", "none"
    errors: list[str] = []
    for source in ("bytes", "gcs_uri"):
        try:
            from google import genai
            from google.genai import types
            client = genai.Client(vertexai=True, project=project, location=os.getenv("GOOGLE_CLOUD_LOCATION") or "global")
            media = (
                types.Part.from_bytes(data=data, mime_type=content_type)
                if source == "bytes"
                else types.Part.from_uri(file_uri=gcs_path, mime_type=content_type)
            )
            response = client.models.generate_content(
                model=model,
                contents=[media, "Return strict JSON with string fields vision_summary and ocr_text. Describe the image accurately and transcribe visible text. Do not add other fields."],
                config=types.GenerateContentConfig(response_mime_type="application/json"),
            )
            parsed = json.loads(response.text or "{}")
            return "ready", str(parsed.get("vision_summary") or "")[:6000], str(parsed.get("ocr_text") or "")[:6000], "", source
        except Exception as exc:
            logger.warning("Attachment vision failed source=%s: %s", source, exc)
            errors.append(f"{source}: {type(exc).__name__}")
    return "failed", "", "", "; ".join(errors)[:300], "none"


def _content_sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _pdf_page_count(data: bytes) -> int:
    from pypdf import PdfReader
    return max(1, len(PdfReader(io.BytesIO(data)).pages))


def _estimate_document_tokens(data: bytes, content_type: str, extracted_chars: int) -> int:
    """Native-read cost estimate: PDFs are billed per page; text-like per char."""
    if content_type == "application/pdf":
        try:
            return _pdf_page_count(data) * TOKENS_PER_PDF_PAGE
        except Exception:
            # Signature already validated; fall back to a size-derived guess.
            return max(TOKENS_PER_PDF_PAGE, (len(data) // 3_000) * TOKENS_PER_PDF_PAGE)
    return max(1, math.ceil(extracted_chars / CHARS_PER_TOKEN))


def _find_reusable_duplicate(batch_id: str, chat_id: str, lecturer_id: str, digest: str) -> ChatAttachment | None:
    """Return an existing unsent, unexpired attachment with the same content hash."""
    docs = (
        _chat_ref(batch_id, chat_id)
        .collection(ATTACHMENTS_SUBCOLLECTION)
        .where("content_sha256", "==", digest)
        .limit(5)
        .stream()
    )
    for doc in docs:
        data = doc.to_dict() or {}
        if (
            data.get("lecturer_id") == lecturer_id
            and data.get("scope") == "chat"
            and not data.get("message_id")
            and not _is_expired(data)
        ):
            return attachment_to_model(doc.id, data)
    return None


def _reserve_chat_storage(batch_id: str, chat_id: str, size: int) -> None:
    """Transactionally reserve quota via a counter on the chat doc (no N-read scan)."""
    db = get_firestore()
    chat_ref = _chat_ref(batch_id, chat_id)
    transaction = db.transaction()

    @firestore.transactional
    def _commit(txn):
        snap = chat_ref.get(transaction=txn)
        current = int((snap.to_dict() or {}).get("attachment_storage_bytes") or 0)
        if current + size > MAX_CHAT_ATTACHMENT_BYTES:
            raise AttachmentValidationError("This chat has reached its 100 MB attachment quota.")
        txn.update(chat_ref, {"attachment_storage_bytes": current + size})

    _commit(transaction)


def _release_chat_storage(batch_id: str, chat_id: str, size: int) -> None:
    try:
        _chat_ref(batch_id, chat_id).update({"attachment_storage_bytes": firestore.Increment(-size)})
    except Exception:
        logger.warning("Failed to release %s bytes of attachment quota chat_id=%s", size, chat_id)


def create_chat_attachment(
    *, batch_id: str, chat_id: str, lecturer_id: str, file_name: str,
    file_title: str, content_type: str, data: bytes,
) -> ChatAttachment:
    """Fast path: validate, store, and record with status=processing.

    Heavy derivatives (extraction, thumbnail, vision, token estimate) run in
    process_chat_attachment as a background task.
    """
    clean_name, normalized_type = validate_attachment(file_name, content_type, data)
    digest = _content_sha256(data)
    duplicate = _find_reusable_duplicate(batch_id, chat_id, lecturer_id, digest)
    if duplicate is not None:
        return duplicate
    _reserve_chat_storage(batch_id, chat_id, len(data))
    attachment_id = str(uuid.uuid4())
    kind = "image" if normalized_type in IMAGE_CONTENT_TYPES else "document"
    try:
        blob_path = chat_attachment_blob_path(lecturer_id, batch_id, chat_id, attachment_id, clean_name)
        gcs_path = upload_bytes(blob_path, data, normalized_type)
    except Exception:
        _release_chat_storage(batch_id, chat_id, len(data))
        raise

    expires_at = datetime.now(timezone.utc) + timedelta(days=get_chat_attachment_retention_days())
    payload = {
        "attachment_id": attachment_id, "batch_id": batch_id, "chat_id": chat_id,
        "message_id": None, "lecturer_id": lecturer_id, "file_name": clean_name,
        "file_title": (file_title or clean_name).strip() or clean_name,
        "content_type": normalized_type, "size_bytes": len(data), "gcs_path": gcs_path,
        "thumbnail_gcs_path": None, "scope": "chat", "attachment_kind": kind,
        "status": "processing", "content_sha256": digest, "token_estimate": 0,
        "parse_status": "pending" if kind == "document" else "skipped",
        "vision_status": "pending" if kind == "image" else "skipped",
        "extracted_text_path": None, "extracted_text_preview": "",
        "vision_summary": "", "ocr_text": "", "expires_at": expires_at,
        "vision_error": "", "vision_source": "none",
        "promoted_file_id": None, "created_at": SERVER_TIMESTAMP, "updated_at": SERVER_TIMESTAMP,
        # Native-first path: the chunk/embed/OCR pipeline is never scheduled.
        "rag_status": "skipped", "chunk_status": "skipped", "embedding_status": "skipped",
        "semantic_search_ready": False, "chunk_count": 0, "indexed_chars": 0,
        "ocr_status": "skipped" if normalized_type == "application/pdf" else "not_needed",
        "rag_error": "", "rag_updated_at": SERVER_TIMESTAMP,
    }
    ref = attachment_ref(batch_id, chat_id, attachment_id)
    ref.set(payload)
    snap = ref.get()
    return attachment_to_model(attachment_id, snap.to_dict() or payload)


def process_chat_attachment(batch_id: str, chat_id: str, attachment_id: str, data: bytes) -> None:
    """Background derivative processing; transitions status processing -> ready|failed.

    Documents: text extraction (preview + full-text blob for non-native types)
    and token estimate. Images: thumbnail + vision summary (vision failure does
    not fail the attachment; the agent analyzes image bytes at query time).
    """
    ref = attachment_ref(batch_id, chat_id, attachment_id)
    snap = ref.get()
    doc = snap.to_dict() or {}
    if not snap.exists:
        logger.warning("process_chat_attachment: doc missing attachment_id=%s", attachment_id)
        return
    kind = str(doc.get("attachment_kind") or "document")
    content_type = str(doc.get("content_type") or "application/octet-stream")
    lecturer_id = str(doc.get("lecturer_id") or "")
    updates: dict[str, Any] = {"updated_at": SERVER_TIMESTAMP}
    try:
        if kind == "image":
            try:
                thumb_path = chat_attachment_thumbnail_path(lecturer_id, batch_id, chat_id, attachment_id)
                updates["thumbnail_gcs_path"] = upload_bytes(thumb_path, _thumbnail_bytes(data), "image/jpeg")
            except Exception as exc:
                logger.warning("Thumbnail generation failed attachment_id=%s: %s", attachment_id, exc)
            vision_status, vision_summary, ocr_text, vision_error, vision_source = _vision_context(
                data, str(doc.get("gcs_path") or ""), content_type
            )
            updates.update({
                "vision_status": vision_status, "vision_summary": vision_summary,
                "ocr_text": ocr_text, "vision_error": vision_error,
                "vision_source": vision_source,
                "token_estimate": IMAGE_TOKEN_ESTIMATE, "status": "ready",
            })
        else:
            full_text = extract_document(data, content_type, MAX_FULL_EXTRACT_CHARS).text.strip()
            updates["extracted_text_preview"] = full_text[:MAX_EXTRACTED_PREVIEW_CHARS]
            if content_type not in NATIVE_READABLE_CONTENT_TYPES:
                # DOCX/PPTX: Gemini cannot read these natively; store the full
                # extracted text as the native-read artifact.
                text_blob_path = chat_attachment_blob_path(
                    lecturer_id, batch_id, chat_id, attachment_id, "extracted_text.txt"
                )
                updates["extracted_text_path"] = upload_bytes(
                    text_blob_path, full_text.encode("utf-8"), "text/plain"
                )
            updates.update({
                "token_estimate": _estimate_document_tokens(data, content_type, len(full_text)),
                "parse_status": "ready", "status": "ready",
            })
    except Exception as exc:
        logger.warning("Attachment processing failed attachment_id=%s: %s", attachment_id, exc)
        updates.update({
            "status": "failed",
            "parse_status": "failed" if kind == "document" else str(doc.get("parse_status") or "skipped"),
        })
    ref.update(updates)


def get_chat_attachment(batch_id: str, chat_id: str, attachment_id: str, lecturer_id: str) -> ChatAttachment | None:
    snap = attachment_ref(batch_id, chat_id, attachment_id).get()
    if not snap.exists:
        return None
    data = snap.to_dict() or {}
    if data.get("lecturer_id") != lecturer_id or data.get("batch_id") != batch_id or data.get("chat_id") != chat_id:
        return None
    return attachment_to_model(snap.id, data)


def _is_expired(data: dict[str, Any]) -> bool:
    expires = data.get("expires_at")
    if isinstance(expires, str):
        try: expires = datetime.fromisoformat(expires.replace("Z", "+00:00"))
        except ValueError: return True
    return bool(expires and expires <= datetime.now(timezone.utc))


def _agent_safe_attachment(doc_id: str, data: dict[str, Any], include_context: bool = True) -> dict[str, Any]:
    result = {
        "attachment_id": doc_id, "message_id": str(data.get("message_id") or ""),
        "file_name": str(data.get("file_name") or ""),
        "file_title": str(data.get("file_title") or data.get("file_name") or ""),
        "content_type": str(data.get("content_type") or ""),
        "size_bytes": int(data.get("size_bytes") or 0),
        "attachment_kind": str(data.get("attachment_kind") or "other"),
        "status": str(data.get("status") or "processing"),
        "token_estimate": int(data.get("token_estimate") or 0),
        "preview_only": True,
        "parse_status": str(data.get("parse_status") or "skipped"),
        "vision_status": str(data.get("vision_status") or "skipped"),
        "vision_source": str(data.get("vision_source") or "none"),
        "thumbnail_available": bool(data.get("thumbnail_gcs_path")),
        "rag_status": str(data.get("rag_status") or "skipped"),
        "chunk_status": str(data.get("chunk_status") or "skipped"),
        "embedding_status": str(data.get("embedding_status") or "skipped"),
        "semantic_search_ready": bool(data.get("semantic_search_ready", False)),
        "chunk_count": int(data.get("chunk_count") or 0),
        "indexed_chars": int(data.get("indexed_chars") or 0),
        "ocr_status": str(data.get("ocr_status") or "not_needed"),
        "rag_updated_at": _iso(data.get("rag_updated_at")),
        "created_at": _iso(data.get("created_at")), "expires_at": _iso(data.get("expires_at")),
    }
    if include_context:
        result.update({
            "extracted_text_preview": str(data.get("extracted_text_preview") or "")[:MAX_EXTRACTED_PREVIEW_CHARS],
            "vision_summary": str(data.get("vision_summary") or "")[:6000],
            "ocr_text": str(data.get("ocr_text") or "")[:6000],
        })
    return result


def _owned_visible_chat(batch_id: str, chat_id: str, lecturer_id: str) -> bool:
    snap = _chat_ref(batch_id, chat_id).get(); data = snap.to_dict() or {}
    return bool(snap.exists and data.get("lecturer_id") == lecturer_id and not data.get("hidden", False))


def list_live_attachment_docs(batch_id: str, chat_id: str, lecturer_id: str, limit: int = 50) -> list[dict[str, Any]]:
    """Raw Firestore docs (incl. gcs_path) for the gateway's trusted manifest.

    Internal only — never expose these dicts to HTTP responses or the model.
    """
    if not _owned_visible_chat(batch_id, chat_id, lecturer_id):
        return []
    safe_limit = max(1, min(int(limit), 50))
    results: list[dict[str, Any]] = []
    docs = _chat_ref(batch_id, chat_id).collection(ATTACHMENTS_SUBCOLLECTION).order_by("created_at", direction="DESCENDING").limit(100).stream()
    for doc in docs:
        data = doc.to_dict() or {}
        if data.get("lecturer_id") != lecturer_id or data.get("scope") != "chat" or _is_expired(data):
            continue
        data["attachment_id"] = doc.id
        results.append(data)
        if len(results) >= safe_limit:
            break
    return results


def list_chat_attachments_for_agent(batch_id: str, chat_id: str, lecturer_id: str, limit: int = 50) -> list[dict[str, Any]]:
    if not _owned_visible_chat(batch_id, chat_id, lecturer_id): return []
    safe_limit = max(1, min(int(limit), 50)); results: list[dict[str, Any]] = []
    docs = _chat_ref(batch_id, chat_id).collection(ATTACHMENTS_SUBCOLLECTION).order_by("created_at", direction="DESCENDING").limit(100).stream()
    for doc in docs:
        data = doc.to_dict() or {}
        if data.get("lecturer_id") != lecturer_id or data.get("scope") != "chat" or _is_expired(data): continue
        results.append(_agent_safe_attachment(doc.id, data, include_context=False))
        if len(results) >= safe_limit: break
    return results


def get_chat_attachment_for_agent(batch_id: str, chat_id: str, lecturer_id: str, attachment_id: str) -> dict[str, Any] | None:
    if not _owned_visible_chat(batch_id, chat_id, lecturer_id): return None
    snap = attachment_ref(batch_id, chat_id, attachment_id).get(); data = snap.to_dict() or {}
    if not snap.exists or data.get("lecturer_id") != lecturer_id or data.get("scope") != "chat" or _is_expired(data): return None
    return _agent_safe_attachment(snap.id, data)


def search_chat_attachments_for_agent(batch_id: str, chat_id: str, lecturer_id: str, query: str, limit: int = 10) -> dict[str, Any]:
    from services.chat_vector_search import search_chat_attachment_chunks
    rag = search_chat_attachment_chunks(batch_id, chat_id, lecturer_id, query, max_results=limit)
    if rag.get("hits"):
        return rag
    if not _owned_visible_chat(batch_id, chat_id, lecturer_id):
        return {"status": "empty", "query": query[:500], "hits": []}
    terms = set(re.findall(r"[\w-]{2,}", query.lower())); scored: list[tuple[int, dict[str, Any]]] = []
    docs = _chat_ref(batch_id, chat_id).collection(ATTACHMENTS_SUBCOLLECTION).order_by("created_at", direction="DESCENDING").limit(100).stream()
    for doc in docs:
        data = doc.to_dict() or {}
        if data.get("lecturer_id") != lecturer_id or data.get("scope") != "chat" or _is_expired(data): continue
        item = _agent_safe_attachment(doc.id, data, include_context=True)
        fields = {
            "file_name": item["file_name"], "file_title": item["file_title"],
            "extracted_text_preview": item.get("extracted_text_preview", ""),
            "vision_summary": item.get("vision_summary", ""), "ocr_text": item.get("ocr_text", ""),
        }
        matched = [name for name, value in fields.items() if terms & set(re.findall(r"[\w-]{2,}", str(value).lower()))]
        score = sum(len(terms & set(re.findall(r"[\w-]{2,}", str(value).lower()))) for value in fields.values())
        if terms and score == 0: continue
        source = next((str(fields[name]) for name in ("extracted_text_preview", "vision_summary", "ocr_text", "file_title", "file_name") if name in matched and fields[name]), item["file_title"])
        scored.append((score, {key: item[key] for key in ("attachment_id", "message_id", "file_name", "file_title", "attachment_kind", "content_type", "created_at", "expires_at")} | {"matched_fields": matched, "snippet": source[:1500], "score": score}))
    hits = [item for _, item in sorted(scored, key=lambda entry: (-entry[0], str(entry[1].get("created_at") or "")), reverse=False)[:max(1, min(limit, 10))]]
    return {"status": "success" if hits else "empty", "query": query[:500], "hits": hits, "retrieval_mode": "lexical"}


def get_attachment_url(batch_id: str, chat_id: str, attachment_id: str, lecturer_id: str, thumbnail: bool) -> str | None:
    snap = attachment_ref(batch_id, chat_id, attachment_id).get(); data = snap.to_dict() or {}
    if not snap.exists or data.get("lecturer_id") != lecturer_id or data.get("batch_id") != batch_id or data.get("chat_id") != chat_id or _is_expired(data):
        return None
    path = data.get("thumbnail_gcs_path") if thumbnail else data.get("gcs_path")
    return signed_read_url(path) if path else None


def delete_attachment_record(batch_id: str, chat_id: str, attachment_id: str, lecturer_id: str | None = None, require_unsent: bool = False) -> str:
    """Delete GCS objects then the Firestore record; return deleted/not_found/sent/storage_failed."""
    ref = attachment_ref(batch_id, chat_id, attachment_id)
    snap = ref.get()
    if not snap.exists:
        return "not_found"
    data = snap.to_dict() or {}
    if lecturer_id and data.get("lecturer_id") != lecturer_id:
        return "not_found"
    if require_unsent and data.get("message_id"):
        return "sent"
    paths = [data.get("gcs_path"), data.get("thumbnail_gcs_path"), data.get("extracted_text_path")]
    if not all(delete_blob(str(path)) for path in paths if path):
        return "storage_failed"
    from services.chat_file_rag_service import delete_attachment_chunks
    if not delete_attachment_chunks(batch_id, chat_id, attachment_id):
        return "storage_failed"
    ref.delete()
    _release_chat_storage(batch_id, chat_id, int(data.get("size_bytes") or 0))
    return "deleted"


def delete_all_chat_attachments(batch_id: str, chat_id: str) -> None:
    for doc in _chat_ref(batch_id, chat_id).collection(ATTACHMENTS_SUBCOLLECTION).stream():
        delete_attachment_record(batch_id, chat_id, doc.id)


def cleanup_expired_attachments(limit: int = 100) -> int:
    db = get_firestore()
    now = datetime.now(timezone.utc)
    retention_days = get_chat_attachment_retention_days()
    cutoff = now - timedelta(days=retention_days)
    docs = db.collection_group(ATTACHMENTS_SUBCOLLECTION).where("expires_at", "<=", now).limit(limit).stream()
    cleaned = 0
    for doc in docs:
        data = doc.to_dict() or {}
        if data.get("scope") != "chat":
            continue
        chat_snap = doc.reference.parent.parent.get()
        chat_data = chat_snap.to_dict() or {}
        updated = chat_data.get("updated_at")
        if updated and updated > cutoff:
            expires_at = updated + timedelta(days=retention_days)
            from services.chat_file_rag_service import update_attachment_chunks_expiry
            try:
                update_attachment_chunks_expiry(str(data.get("batch_id") or ""), str(data.get("chat_id") or ""), doc.id, expires_at)
            except Exception:
                logger.exception("Failed to extend attachment chunk retention attachment_id=%s", doc.id)
                continue
            doc.reference.update({"expires_at": expires_at, "updated_at": SERVER_TIMESTAMP})
            continue
        if delete_attachment_record(str(data.get("batch_id") or ""), str(data.get("chat_id") or ""), doc.id) in {"deleted", "not_found"}:
            cleaned += 1
    return cleaned
