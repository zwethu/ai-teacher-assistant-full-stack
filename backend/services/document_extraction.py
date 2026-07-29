"""Bounded reusable extraction for chat previews and course overlays."""
from __future__ import annotations
import csv, io, json
from dataclasses import dataclass

@dataclass(frozen=True)
class ExtractedSegment:
    text: str
    page_number: int | None = None

@dataclass(frozen=True)
class ExtractionResult:
    segments: list[ExtractedSegment]
    text: str
    truncated: bool

def extract_document(data: bytes, content_type: str, max_chars: int) -> ExtractionResult:
    segments: list[ExtractedSegment] = []
    consumed = 0
    truncated = False
    def add(text: str, page: int | None = None) -> bool:
        nonlocal consumed, truncated
        clean = str(text or "").strip()
        if not clean: return True
        remaining = max_chars - consumed
        if remaining <= 0: truncated = True; return False
        if len(clean) > remaining: clean = clean[:remaining]; truncated = True
        segments.append(ExtractedSegment(clean, page)); consumed += len(clean) + 2
        return not truncated
    if content_type == "application/pdf":
        from pypdf import PdfReader
        for index, page in enumerate(PdfReader(io.BytesIO(data)).pages, 1):
            if not add(page.extract_text() or "", index): break
    elif "wordprocessingml" in content_type:
        from docx import Document
        for paragraph in Document(io.BytesIO(data)).paragraphs:
            if not add(paragraph.text): break
    elif "presentationml" in content_type:
        from pptx import Presentation
        for index, slide in enumerate(Presentation(io.BytesIO(data)).slides, 1):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and str(shape.text or "").strip())
            if not add(text, index): break
    else:
        decoded = data.decode("utf-8-sig", errors="strict")
        if content_type == "text/csv":
            for row in csv.reader(io.StringIO(decoded)):
                if not add(" | ".join(cell.strip() for cell in row)): break
        elif content_type == "application/json": add(json.dumps(json.loads(decoded), ensure_ascii=False, indent=2))
        else: add(decoded)
    text = "\n\n".join(item.text for item in segments)[:max_chars]
    return ExtractionResult(segments, text, truncated)

def chunk_extraction(result: ExtractionResult, target_chars: int = 4000, overlap_chars: int = 300, max_chunks: int = 100) -> tuple[list[ExtractedSegment], bool]:
    chunks: list[ExtractedSegment] = []; truncated = result.truncated
    for segment_index, segment in enumerate(result.segments):
        text = segment.text.strip(); start = 0
        while start < len(text):
            end = min(start + target_chars, len(text))
            if end < len(text):
                split = text.rfind("\n", start + 3000, min(start + 5000, len(text)))
                if split <= start: split = text.rfind(" ", start + 3000, min(start + 5000, len(text)))
                if split > start: end = split
            chunks.append(ExtractedSegment(text[start:end].strip(), segment.page_number))
            if len(chunks) >= max_chunks:
                return chunks, truncated or end < len(text) or segment_index < len(result.segments) - 1
            if end >= len(text): break
            start = max(end - overlap_chars, start + 1)
    return chunks, truncated
